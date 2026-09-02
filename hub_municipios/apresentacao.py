"""
Gerador da apresentação comercial da carteira — .pptx editável.

Público: **prefeitura / poder concedente**. A peça responde à pergunta que a cidade
faz, não à que o investidor faz: "a CIP que já arrecadamos banca a modernização da
nossa iluminação?". Por isso o slide fala de economia de energia, qualidade do serviço
e uso da receita existente — e **não** expõe margem de concessionário, TIR nem o
ranking comercial da carteira, que são conversa de outra mesa.

Estrutura da peça de cliente:

    Capa → De onde vem esta análise → Possíveis escopos
      → [ município: panorama | pré-viabilidade | impactos | a conta fecha? |
          soluções digitais ] × N
      → Metodologia e fontes

## Por que vetor, e não imagem

Tudo é desenhado com formas nativas do PowerPoint — inclusive o mapa, construído com
`build_freeform` a partir da malha do IBGE. Não há PNG em lugar nenhum. Isso significa
que quem receber o arquivo pode mover, recolorir e editar qualquer elemento, o texto
continua selecionável e a impressão sai nítida em qualquer tamanho. Também evita a
dependência do `kaleido` (exportador de imagem do Plotly), que não está instalado e
pesa ~100 MB no deploy.

O mapa usa o contorno da UF como um polígono só (`malhas.carregar_contorno_uf`), e não
os 800 polígonos municipais: empilhar 800 formas por slide inviabilizaria o arquivo.

## Honestidade dos números

O que sai no slide é o que o Hub apurou, com a mesma disciplina do resto do módulo:

  - Número ausente vira "—", nunca zero nem estimativa silenciosa.
  - Parque estimado por população (sem BDGD) é rotulado como estimativa no rodapé.
  - Município com declaração implausível ou potência fora da faixa física **não entra**
    na apresentação: a UI filtra antes, e aqui o rodapé denuncia se algo passar.
  - A cobertura da CIP sobre a conta de energia só aparece onde o município declarou a
    função 25 ao SICONFI — ver `despesas.py`, que mede a ausência em 45% dos casos.

## Padrão visual

O layout replica o deck "Estruturação da PPP de Cidade Inteligente — Itabira"
(05/08/2026), adotado pelo usuário em 02/09/2026 como padrão da peça comercial:
faixa-título arredondada em degradê azul, fundo off-white com arcos tênues, par
verde/vermelho para arrecadação × despesa e rodapé da marca. Paleta extraída dos
próprios slides, não estimada de memória.
"""

from __future__ import annotations

import io
from dataclasses import dataclass
from datetime import datetime
from typing import Any, Iterable, Optional, Sequence

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import malhas, ranking


# ── Identidade da peça comercial ─────────────────────────────────────────────
# EXCEÇÃO DELIBERADA à diretriz de 30/08/2026 (o portal não gera nada com a marca
# Houer), autorizada pelo usuário em 02/09/2026: esta é peça COMERCIAL, apresentada
# pela consultoria ao município, e não um arquivo do portal. A exceção vale só para o
# .pptx — planilhas, relatórios e a UI continuam como Plataforma IP.
#
# As cores foram amostradas pixel a pixel do deck de Itabira, não estimadas: o azul da
# faixa é um degradê (#003264 → #00508E), o verde da arrecadação é saturado (#00FF29) e
# o vermelho da despesa é #FF3131. Trocar por "azul corporativo genérico" descaracteriza
# a peça ao lado das que a equipe já entregou.
NAVY = RGBColor(0x00, 0x32, 0x64)          # texto institucional e ponta escura do degradê
NAVY_ESCURO = RGBColor(0x00, 0x2D, 0x56)   # painel da capa
AZUL = RGBColor(0x00, 0x50, 0x8E)          # ponta clara do degradê da faixa
CIANO = RGBColor(0x00, 0xB0, 0xF0)         # chip de data, detalhes
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
TINTA = RGBColor(0x0D, 0x29, 0x53)         # número grande
TINTA_FRACA = RGBColor(0x64, 0x74, 0x8B)
FUNDO = RGBColor(0xF6, 0xF6, 0xF6)         # off-white de todos os slides internos
CINZA_CLARO = RGBColor(0xE2, 0xE8, 0xF0)
CINZA_ARCO = RGBColor(0xEC, 0xEC, 0xEC)    # arcos tênues do fundo
CINZA_MAPA = RGBColor(0xCB, 0xD5, 0xE1)
MARCA_FRACA = RGBColor(0xD6, 0xD6, 0xD6)   # rodapé da marca, propositalmente discreto
VERDE = RGBColor(0x00, 0xFF, 0x29)         # arrecadação
VERDE_ESCURO = RGBColor(0x0F, 0x9D, 0x3A)  # verde legível sobre fundo claro em texto miúdo
VERMELHO = RGBColor(0xFF, 0x31, 0x31)      # despesa
AMBAR = RGBColor(0xD9, 0x77, 0x06)
ROXO = RGBColor(0x6B, 0x2F, 0xA0)          # pílula de Governança (p05)
VERDE_ESCOPO = RGBColor(0x4C, 0x93, 0x2B)  # pílulas verdes dos escopos (p05)

# Assinatura da marca no rodapé de todo slide interno, como no deck de Itabira.
MARCA_NOME = "HOUER"
MARCA_ASSINATURA = "HOUER  ·  IMPACTANDO GERAÇÕES"
MARCA_SITE = "HOUER.COM.BR"

# Abaixo deste ganho de eficiência a tese energética não se sustenta, e o slide troca
# de argumento. Em MG isso é regra, não exceção: 75,6% do parque estadual já é LED, e
# município com potência média abaixo da referência não tem economia a capturar.
LIMIAR_ECONOMIA_RELEVANTE = 0.05

FONTE = "Inter"          # cai para a fonte padrão do sistema se não existir
FONTE_ALT = "Calibri"

LARGURA = Inches(13.333)
ALTURA = Inches(7.5)
MARGEM = Inches(0.6)


# ── Soluções digitais custeáveis pela sobra da CIP ───────────────────────────
# A CIP é vinculada ao serviço de iluminação pública (art. 149-A da Constituição), e o
# que sobra dela pode custear o que se apoia na rede de IP — não qualquer despesa da
# prefeitura. Por isso a lista abaixo é de soluções que compartilham poste, energia e
# conectividade com a iluminação, e não um catálogo genérico de "cidade inteligente".
# Cada item vem SEM custo por decisão do usuário em 01/09/2026: o slide apresenta o
# que é possível, e o dimensionamento fica para a proposta.
SOLUCOES_DIGITAIS = [
    ("Telegestão", "Comando e medição ponto a ponto: dimerização por horário, "
                   "detecção automática de falha e leitura de consumo real."),
    ("Videomonitoramento", "Câmeras no poste, integradas ao centro de controle "
                           "municipal e à segurança pública."),
    ("Wi-Fi público", "Pontos de acesso em praças e equipamentos públicos, usando a "
                      "rede elétrica e a fibra da própria concessão."),
    ("Sensoriamento urbano", "Qualidade do ar, ruído, alagamento e contagem de "
                             "tráfego, com dados abertos ao município."),
    ("Botão de pânico e SOS", "Acionamento em ponto fixo com vídeo associado, "
                              "integrado ao atendimento de emergência."),
    ("Painel de gestão", "Indicadores de operação, prazo de atendimento e consumo "
                         "em tempo real, para fiscalizar o contrato."),
]


@dataclass
class Bloco:
    """Um número em destaque no slide."""

    rotulo: str
    valor: str
    apoio: str = ""
    cor: Optional[RGBColor] = None


# ── Helpers de desenho ───────────────────────────────────────────────────────
def _fundo(slide, cor: RGBColor) -> None:
    fundo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, LARGURA, ALTURA)
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = cor
    fundo.line.fill.background()
    fundo.shadow.inherit = False
    slide.shapes._spTree.remove(fundo._element)
    slide.shapes._spTree.insert(2, fundo._element)


def _degrade(forma, inicio: RGBColor, fim: RGBColor, angulo: float = 0.0) -> None:
    """
    Preenche a forma com um degradê linear de dois passos.

    A faixa-título do padrão Itabira é um degradê, não uma cor chapada: sem ele a peça
    fica visivelmente diferente das que a equipe já entregou. `angulo=0` corre da
    esquerda para a direita, que é o sentido usado no deck.
    """
    forma.fill.gradient()
    passos = forma.fill.gradient_stops
    passos[0].color.rgb = inicio
    passos[0].position = 0.0
    passos[1].color.rgb = fim
    passos[1].position = 1.0
    # `gradient()` do python-pptx cria três paradas; a do meio deixaria um vinco.
    for extra in list(passos)[2:]:
        extra.color.rgb = fim
        extra.position = 1.0
    forma.fill.gradient_angle = angulo
    forma.line.fill.background()
    forma.shadow.inherit = False


def _texto(slide, x, y, w, h, texto: str, *, tamanho=14, cor=TINTA, negrito=False,
           alinhamento=PP_ALIGN.LEFT, fonte=FONTE, espacamento=1.0,
           ancora=None, entreletras: Optional[int] = None):
    caixa = slide.shapes.add_textbox(x, y, w, h)
    quadro = caixa.text_frame
    quadro.word_wrap = True
    quadro.margin_left = quadro.margin_right = 0
    quadro.margin_top = quadro.margin_bottom = 0
    if ancora is not None:
        quadro.vertical_anchor = ancora
    paragrafo = quadro.paragraphs[0]
    paragrafo.alignment = alinhamento
    paragrafo.line_spacing = espacamento
    corrida = paragrafo.add_run()
    corrida.text = texto
    corrida.font.size = Pt(tamanho)
    corrida.font.bold = negrito
    corrida.font.color.rgb = cor
    corrida.font.name = fonte
    if entreletras:
        # `spc` é centésimos de ponto e não tem API no python-pptx. É o que dá o ar
        # institucional da assinatura da marca no rodapé do deck.
        corrida.font._rPr.set("spc", str(int(entreletras)))
    return caixa


def _retangulo(slide, x, y, w, h, cor, *, borda=None, arredondado=False):
    forma = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if arredondado else MSO_SHAPE.RECTANGLE, x, y, w, h)
    forma.fill.solid()
    forma.fill.fore_color.rgb = cor
    if borda is None:
        forma.line.fill.background()
    else:
        forma.line.color.rgb = borda
        forma.line.width = Pt(1)
    forma.shadow.inherit = False
    if arredondado:
        forma.adjustments[0] = 0.06
    return forma


def _cartao(slide, x, y, w, h, bloco: Bloco) -> None:
    """Um cartão de indicador: rótulo pequeno, número grande, apoio menor."""
    _retangulo(slide, x, y, w, h, BRANCO, borda=CINZA_CLARO, arredondado=True)
    _retangulo(slide, x, y, Inches(0.05), h, bloco.cor or CIANO, arredondado=False)
    pad = Inches(0.22)
    _texto(slide, x + pad, y + Inches(0.16), w - 2 * pad, Inches(0.3),
           bloco.rotulo.upper(), tamanho=9, cor=TINTA_FRACA, negrito=True)
    _texto(slide, x + pad, y + Inches(0.46), w - 2 * pad, Inches(0.5),
           bloco.valor, tamanho=25, cor=bloco.cor or NAVY, negrito=True)
    if bloco.apoio:
        _texto(slide, x + pad, y + h - Inches(0.42), w - 2 * pad, Inches(0.3),
               bloco.apoio, tamanho=9, cor=TINTA_FRACA)


def _arcos_de_fundo(slide) -> None:
    """
    Arcos tênues do fundo, marca d'água do padrão Itabira.

    São círculos gigantes com centro fora do slide: só a curva entra no quadro, que é
    exatamente o efeito do original. Sem preenchimento, para não pesar o arquivo nem
    interferir na leitura do texto por cima.
    """
    for esquerda, topo, diametro in (
        (Inches(-3.2), Inches(-4.6), Inches(12.5)),
        (Inches(-1.4), Inches(-3.4), Inches(11.0)),
        (Inches(7.4), Inches(-2.2), Inches(12.0)),
    ):
        arco = slide.shapes.add_shape(MSO_SHAPE.OVAL, esquerda, topo, diametro, diametro)
        arco.fill.background()
        arco.line.color.rgb = CINZA_ARCO
        arco.line.width = Pt(1.25)
        arco.shadow.inherit = False


def _fundo_padrao(slide) -> None:
    """Off-white com os arcos — o fundo de todo slide interno."""
    _fundo(slide, FUNDO)
    _arcos_de_fundo(slide)


def _faixa_titulo(slide, titulo: str, subtitulo: str = "") -> None:
    """
    Faixa arredondada em degradê, título branco em caixa alta e centralizado.

    É o elemento que mais identifica o padrão Itabira. O subtítulo, que o deck original
    não tem, fica logo abaixo da faixa em azul — dizer de qual município é o slide vale
    mais, numa peça gerada em série, do que a fidelidade absoluta ao modelo.
    """
    faixa = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGEM, Inches(0.16),
        LARGURA - 2 * MARGEM, Inches(1.16))
    faixa.adjustments[0] = 0.5
    _degrade(faixa, NAVY, AZUL)
    faixa.text_frame.text = ""
    _texto(slide, MARGEM + Inches(0.4), Inches(0.16), LARGURA - 2 * MARGEM - Inches(0.8),
           Inches(1.16), titulo.upper(), tamanho=30, cor=BRANCO, negrito=True,
           alinhamento=PP_ALIGN.CENTER, ancora=MSO_ANCHOR.MIDDLE, espacamento=0.95)
    if subtitulo:
        _texto(slide, MARGEM, Inches(1.42), LARGURA - 2 * MARGEM, Inches(0.3),
               subtitulo, tamanho=12, cor=AZUL)


def _rodape_marca(slide) -> None:
    """Assinatura da marca nos dois cantos inferiores, como no deck de Itabira."""
    _texto(slide, MARGEM, ALTURA - Inches(0.42), Inches(5), Inches(0.26),
           MARCA_ASSINATURA, tamanho=8, cor=MARCA_FRACA, negrito=True, entreletras=120)
    _texto(slide, LARGURA - MARGEM - Inches(3), ALTURA - Inches(0.42), Inches(3),
           Inches(0.26), MARCA_SITE, tamanho=8, cor=MARCA_FRACA, negrito=True,
           alinhamento=PP_ALIGN.RIGHT, entreletras=120)


def _rodape(slide, texto: str) -> None:
    """Nota técnica acima da assinatura da marca — fontes, ressalvas e origem do dado."""
    _texto(slide, MARGEM, ALTURA - Inches(0.86), LARGURA - 2 * MARGEM, Inches(0.4),
           texto, tamanho=8, cor=TINTA_FRACA, espacamento=1.15)
    _rodape_marca(slide)


def _marca_houer(slide, x, y, w, h) -> None:
    """
    Assinatura da consultoria na capa: caixa branca com borda navy, como no original.

    O logotipo é desenhado como texto, não como imagem: manter a peça 100% vetorial vale
    mais do que reproduzir o símbolo, e o arquivo segue editável para quem quiser colar
    o logotipo oficial por cima.
    """
    _retangulo(slide, x, y, w, h, BRANCO, borda=NAVY, arredondado=True)
    _texto(slide, x, y + Inches(0.16), w, Inches(0.42), MARCA_NOME, tamanho=25,
           cor=NAVY, negrito=True, alinhamento=PP_ALIGN.CENTER, entreletras=60)
    _texto(slide, x, y + Inches(0.6), w, Inches(0.24), "CONSULTORIA", tamanho=9,
           cor=NAVY, alinhamento=PP_ALIGN.CENTER, entreletras=260)


# ── Mapa vetorial ────────────────────────────────────────────────────────────
def _aneis(geometria: dict) -> list[list[tuple[float, float]]]:
    """Extrai os anéis externos de um Polygon ou MultiPolygon GeoJSON."""
    if not geometria:
        return []
    tipo = geometria.get("type")
    coordenadas = geometria.get("coordinates") or []
    if tipo == "Polygon":
        return [[(float(p[0]), float(p[1])) for p in anel] for anel in coordenadas[:1]]
    if tipo == "MultiPolygon":
        return [[(float(p[0]), float(p[1])) for p in poligono[0]]
                for poligono in coordenadas if poligono]
    return []


def _simplificar(anel: Sequence[tuple[float, float]], maximo: int = 260):
    """
    Reduz vértices por amostragem uniforme.

    A malha de qualidade mínima do IBGE já vem simplificada (MG tem 372 pontos), mas
    município recortado por rio pode passar de mil. Acima do teto o PowerPoint fica
    lento para editar sem que a diferença visual apareça no tamanho impresso.
    """
    if len(anel) <= maximo:
        return list(anel)
    passo = len(anel) / maximo
    reduzido = [anel[int(i * passo)] for i in range(maximo)]
    if reduzido[0] != reduzido[-1]:
        reduzido.append(reduzido[0])
    return reduzido


def _projetar(anel, caixa, x, y, w, h):
    """
    Projeta lon/lat na área do slide, preservando a proporção geográfica.

    A correção por cos(latitude) evita o achatamento que apareceria ao esticar graus
    de longitude e latitude na mesma escala — sem ela, Minas Gerais sai gorda.
    """
    lon_min, lat_min, lon_max, lat_max = caixa
    lat_media = (lat_min + lat_max) / 2
    import math
    fator_lon = math.cos(math.radians(lat_media))
    largura_geo = max((lon_max - lon_min) * fator_lon, 1e-9)
    altura_geo = max(lat_max - lat_min, 1e-9)
    escala = min(w / largura_geo, h / altura_geo)
    deslocamento_x = x + (w - largura_geo * escala) / 2
    deslocamento_y = y + (h - altura_geo * escala) / 2
    pontos = []
    for lon, lat in anel:
        px = deslocamento_x + (lon - lon_min) * fator_lon * escala
        py = deslocamento_y + (lat_max - lat) * escala      # y do slide cresce para baixo
        pontos.append((Emu(int(px)), Emu(int(py))))
    return pontos


def _caixa_de(aneis: Iterable[Sequence[tuple[float, float]]]):
    lons, lats = [], []
    for anel in aneis:
        for lon, lat in anel:
            lons.append(lon)
            lats.append(lat)
    if not lons:
        return None
    return min(lons), min(lats), max(lons), max(lats)


def _poligono(slide, pontos, cor, borda=None, largura_borda=Pt(0.75)):
    if len(pontos) < 3:
        return None
    construtor = slide.shapes.build_freeform(pontos[0][0], pontos[0][1])
    construtor.add_line_segments(pontos[1:], close=True)
    forma = construtor.convert_to_shape()
    forma.fill.solid()
    forma.fill.fore_color.rgb = cor
    if borda is None:
        forma.line.fill.background()
    else:
        forma.line.color.rgb = borda
        forma.line.width = largura_borda
    forma.shadow.inherit = False
    return forma


def desenhar_mapa(slide, uf: str, cod_ibge: str, x, y, w, h, *,
                  cor_uf: Optional[RGBColor] = None,
                  borda_uf: Optional[RGBColor] = None) -> bool:
    """
    Desenha o contorno da UF com o município destacado, em vetor.

    `cor_uf`/`borda_uf` existem para o mapa da capa, que fica sobre painel navy e
    precisa do contorno claro — no fundo escuro o cinza padrão desaparece.

    Returns:
        True se conseguiu desenhar. False (sem levantar) quando não há malha nem rede —
        o slide segue válido sem o mapa, apenas com o espaço vazio reaproveitado.
    """
    contorno = malhas.carregar_contorno_uf(uf)
    if not contorno or not contorno.get("features"):
        return False
    aneis_uf = _aneis(contorno["features"][0].get("geometry") or {})
    if not aneis_uf:
        return False

    geometria_municipio = malhas.geometria_do_municipio(uf, cod_ibge)
    aneis_municipio = _aneis(geometria_municipio) if geometria_municipio else []

    caixa = _caixa_de(aneis_uf)
    if not caixa:
        return False

    for anel in aneis_uf:
        _poligono(slide, _projetar(_simplificar(anel), caixa, x, y, w, h),
                  cor_uf or CINZA_MAPA, borda=borda_uf or BRANCO)
    # Anel de destaque ANTES do município: num estado do tamanho de Minas Gerais, um
    # município de 300 km² vira um ponto de 2 mm e o destaque some. O círculo resolve
    # sem falsear a escala — a geometria real continua desenhada por cima dele.
    pontos_municipio = [_projetar(_simplificar(anel, 200), caixa, x, y, w, h)
                        for anel in aneis_municipio]
    if pontos_municipio:
        _anel_de_destaque(slide, [p for pts in pontos_municipio for p in pts])
    for pontos in pontos_municipio:
        _poligono(slide, pontos, CIANO, borda=NAVY, largura_borda=Pt(1.25))
    return True


def _anel_de_destaque(slide, pontos) -> None:
    """Círculo vazado ao redor do município, para o olho achá-lo no mapa do estado."""
    xs = [int(p[0]) for p in pontos]
    ys = [int(p[1]) for p in pontos]
    centro_x, centro_y = (min(xs) + max(xs)) // 2, (min(ys) + max(ys)) // 2
    raio = max(max(xs) - min(xs), max(ys) - min(ys)) // 2 + Inches(0.22)
    circulo = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(centro_x - raio), Emu(centro_y - raio),
        Emu(raio * 2), Emu(raio * 2))
    circulo.fill.background()
    circulo.line.color.rgb = CIANO
    circulo.line.width = Pt(1.75)
    circulo.shadow.inherit = False


# ── Formatação dos indicadores ───────────────────────────────────────────────
def _v(linha, coluna, formato) -> str:
    return ranking.formatar_valor(linha.get(coluna), formato)


def _data_br(valor) -> str:
    """
    ISO → dd/mm/aaaa. Data em formato ISO num slide de prefeitura é ruído: o PNCP
    devolve `2026-01-01`, e é assim que ia parar na tela do cliente.
    """
    texto = str(valor or "")[:10]
    partes = texto.split("-")
    if len(partes) == 3 and len(partes[0]) == 4:
        return f"{partes[2]}/{partes[1]}/{partes[0]}"
    return texto


def _numero(linha, coluna) -> Optional[float]:
    valor = linha.get(coluna)
    try:
        if valor is None or pd.isna(valor):
            return None
        return float(valor)
    except (TypeError, ValueError):
        return None


# ── Slides ───────────────────────────────────────────────────────────────────
def _slide_vazio(apresentacao):
    return apresentacao.slides.add_slide(apresentacao.slide_layouts[6])


def _ondas_capa(slide, x, y, w, h) -> None:
    """
    As faixas onduladas do painel da capa, desenhadas como polígonos.

    Não existe recorte (clipping) no PowerPoint via python-pptx: forma que passa da
    borda do painel invade o slide. Por isso cada faixa é construída já limitada à área
    do painel, amostrando a senoide e travando o x nos limites — em vez de desenhar
    formas grandes e torcer para caberem.
    """
    import math

    passos = 48
    for deslocamento, largura_faixa, cor in (
        (0.30, 0.16, RGBColor(0x00, 0x38, 0x6B)),
        (0.52, 0.14, RGBColor(0x00, 0x42, 0x7C)),
        (0.74, 0.13, RGBColor(0x00, 0x4C, 0x8C)),
    ):
        esquerda, direita = [], []
        for i in range(passos + 1):
            t = i / passos
            onda = 0.07 * math.sin(math.pi * (t * 1.6 + 0.15))
            px_e = x + int(w * min(max(deslocamento + onda, 0.0), 1.0))
            px_d = x + int(w * min(max(deslocamento + onda + largura_faixa, 0.0), 1.0))
            py = y + int(h * t)
            esquerda.append((Emu(px_e), Emu(py)))
            direita.append((Emu(px_d), Emu(py)))
        _poligono(slide, esquerda + list(reversed(direita)), cor)


def _capa(apresentacao, titulo: str, subtitulo: str, municipios: int,
          linha: Optional[dict] = None) -> None:
    """
    Capa no padrão Itabira: painel navy ondulado à esquerda, mapa à direita, chip de
    data e assinatura da consultoria no canto inferior.

    O original traz uma fotografia noturna de via iluminada no quadro da direita. Aqui
    entra o mapa vetorial do município — não por limitação, mas porque foto de banco de
    imagens não diz nada sobre a cidade, e o contorno dela diz.
    """
    slide = _slide_vazio(apresentacao)
    _fundo(slide, FUNDO)
    _arcos_de_fundo(slide)

    painel_x, painel_y = Inches(0.25), Inches(0.25)
    painel_w, painel_h = Inches(7.45), Inches(7.0)
    painel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    painel_x, painel_y, painel_w, painel_h)
    painel.adjustments[0] = 0.03
    _degrade(painel, NAVY_ESCURO, RGBColor(0x00, 0x3D, 0x74))
    painel.text_frame.text = ""
    _ondas_capa(slide, painel_x, painel_y, painel_w, painel_h)

    _texto(slide, painel_x + Inches(0.6), Inches(2.55), painel_w - Inches(1.2),
           Inches(2.4), titulo.upper(), tamanho=40, cor=BRANCO, negrito=True,
           espacamento=1.1)
    _texto(slide, painel_x + Inches(0.6), Inches(5.15), painel_w - Inches(1.2),
           Inches(1.0), subtitulo, tamanho=13, cor=CINZA_CLARO, espacamento=1.3)

    # ── Quadro da direita: mapa ou, sem malha, painel de identificação ──────
    quadro_x, quadro_y = Inches(7.95), Inches(0.25)
    quadro_w, quadro_h = Inches(5.15), Inches(5.55)
    quadro = _retangulo(slide, quadro_x, quadro_y, quadro_w, quadro_h, NAVY_ESCURO,
                        arredondado=True)
    quadro.adjustments[0] = 0.03
    uf = str((linha or {}).get("uf") or "")
    cod = str((linha or {}).get("codigo_municipio") or "")
    desenhou = False
    if uf:
        desenhou = desenhar_mapa(slide, uf, cod, quadro_x + Inches(0.45),
                                 quadro_y + Inches(0.45), quadro_w - Inches(0.9),
                                 quadro_h - Inches(1.35),
                                 cor_uf=RGBColor(0x0B, 0x4C, 0x86),
                                 borda_uf=RGBColor(0x1E, 0x6E, 0xB0))
    if not desenhou:
        # Sem malha (carteira com vários municípios, ou IBGE fora do ar) o quadro
        # ficaria um retângulo navy vazio no meio da capa. A silhueta ocupa o lugar.
        _cidade_vetorial(slide, quadro_x + Inches(0.55), quadro_y + Inches(1.35),
                         quadro_w - Inches(1.1), Inches(3.0), sobre_escuro=True)
    legenda = str((linha or {}).get("municipio") or "").upper()
    if legenda:
        _texto(slide, quadro_x, quadro_y + quadro_h - Inches(0.75), quadro_w,
               Inches(0.4), f"{legenda}  ·  {uf}" if uf else legenda,
               tamanho=13, cor=CIANO, negrito=True, alinhamento=PP_ALIGN.CENTER,
               entreletras=140)

    # Chip de data + assinatura, na mesma linha, como no deck original.
    chip_y = Inches(6.05)
    chip = _retangulo(slide, quadro_x, chip_y, Inches(2.25), Inches(1.2), CIANO)
    chip.line.fill.background()
    rotulo = "município" if municipios == 1 else "municípios"
    _texto(slide, quadro_x, chip_y + Inches(0.28), Inches(2.25), Inches(0.36),
           f"{datetime.now():%d/%m/%Y}", tamanho=13, cor=NAVY, negrito=True,
           alinhamento=PP_ALIGN.CENTER)
    _texto(slide, quadro_x, chip_y + Inches(0.68), Inches(2.25), Inches(0.32),
           f"{municipios} {rotulo}", tamanho=11, cor=NAVY,
           alinhamento=PP_ALIGN.CENTER)
    _marca_houer(slide, quadro_x + Inches(2.45), chip_y, Inches(2.7), Inches(1.2))


def _slide_panorama(apresentacao, linha: dict) -> None:
    """Slide 1: identificação, mapa de localização e os números do município hoje."""
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    municipio = str(linha.get("municipio") or "Município")
    uf = str(linha.get("uf") or "")
    exercicio = linha.get("ano_exercicio")
    _faixa_titulo(slide, f"{municipio}{'/' + uf if uf else ''}",
                  "Panorama da iluminação pública e da arrecadação de CIP")

    topo = Inches(1.9)
    largura_mapa = Inches(4.3)
    desenhou = desenhar_mapa(slide, uf, str(linha.get("codigo_municipio") or ""),
                             MARGEM, topo, largura_mapa, Inches(4.0))
    if desenhou:
        _texto(slide, MARGEM, topo + Inches(4.05), largura_mapa, Inches(0.3),
               f"Localização em {uf}", tamanho=9, cor=TINTA_FRACA,
               alinhamento=PP_ALIGN.CENTER)

    x0 = MARGEM + largura_mapa + Inches(0.5) if desenhou else MARGEM
    largura_util = LARGURA - x0 - MARGEM
    largura_cartao = (largura_util - Inches(0.3)) / 2
    altura_cartao = Inches(1.35)

    origem = str(linha.get("origem_pontos") or "")
    blocos = [
        Bloco("População", _v(linha, "populacao", "numero"), "IBGE"),
        Bloco("Pontos de iluminação", _v(linha, "pontos_ip", "numero"),
              origem or "BDGD/ANEEL"),
        Bloco("Arrecadação de CIP", _v(linha, "cosip_liquida", "reais_compacto"),
              f"exercício {exercicio}" if exercicio else "SICONFI"),
        Bloco("CIP por ponto", _v(linha, "cosip_ponto_mes", "reais"), "por ponto/mês"),
    ]
    for indice, bloco in enumerate(blocos):
        coluna, fila = indice % 2, indice // 2
        _cartao(slide, x0 + coluna * (largura_cartao + Inches(0.3)),
                topo + fila * (altura_cartao + Inches(0.3)),
                largura_cartao, altura_cartao, bloco)

    # Faixa da conta de energia — só quando o município declarou a função 25.
    y_faixa = topo + 2 * (altura_cartao + Inches(0.3))
    cobertura = _numero(linha, "cosip_cobre_energia")
    despesa = _numero(linha, "despesa_energia_declarada")
    if cobertura is not None and despesa is not None:
        cor = VERDE_ESCURO if cobertura >= 1 else AMBAR
        _retangulo(slide, x0, y_faixa, largura_util, Inches(1.05), BRANCO,
                   borda=CINZA_CLARO, arredondado=True)
        _retangulo(slide, x0, y_faixa, Inches(0.05), Inches(1.05), cor)
        verbo = "cobre" if cobertura >= 1 else "cobre apenas"
        _texto(slide, x0 + Inches(0.22), y_faixa + Inches(0.16),
               largura_util - Inches(0.44), Inches(0.3),
               "A CIP FRENTE À CONTA DE ENERGIA DO MUNICÍPIO", tamanho=9,
               cor=TINTA_FRACA, negrito=True)
        _texto(slide, x0 + Inches(0.22), y_faixa + Inches(0.45),
               largura_util - Inches(0.44), Inches(0.5),
               f"A arrecadação de CIP {verbo} "
               f"{ranking.formatar_valor(cobertura, 'percentual')} dos "
               f"{ranking.formatar_valor(despesa, 'reais_compacto')} declarados em "
               "despesa com energia elétrica.",
               tamanho=12, cor=TINTA, espacamento=1.2)

    _rodape(slide, _fonte_rodape(linha))


# ── Glifos dos ícones ────────────────────────────────────────────────────────
# Emblemas circulares com um glifo branco dentro, montados só com formas nativas.
# O deck de Itabira usa ícones de linha de um pacote licenciado; até a equipe passar a
# biblioteca (decisão do usuário em 02/09/2026: seguir com vetor próprio por enquanto),
# estes fazem o papel sem risco de licença e mantêm o arquivo editável.
def _forma_branca(slide, forma_mso, x, y, w, h, cor=BRANCO):
    forma = slide.shapes.add_shape(forma_mso, Emu(int(x)), Emu(int(y)),
                                   Emu(int(w)), Emu(int(h)))
    forma.fill.solid()
    forma.fill.fore_color.rgb = cor
    forma.line.fill.background()
    forma.shadow.inherit = False
    return forma


def _glifo(slide, chave: str, x, y, tamanho) -> None:
    """Desenha o glifo `chave` dentro de um emblema circular navy de lado `tamanho`."""
    emblema = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, tamanho, tamanho)
    emblema.fill.solid()
    emblema.fill.fore_color.rgb = NAVY
    emblema.line.fill.background()
    emblema.shadow.inherit = False

    cx, cy, t = int(x) + int(tamanho) / 2, int(y) + int(tamanho) / 2, int(tamanho)
    if chave == "pessoas":
        # Cabeça e ombros com folga entre as duas formas: como tudo é branco sobre
        # navy, encostar as formas funde as duas num borrão.
        _forma_branca(slide, MSO_SHAPE.OVAL, cx - t * 0.14, cy - t * 0.30,
                      t * 0.28, t * 0.28)
        ombros = _forma_branca(slide, MSO_SHAPE.ROUND_2_SAME_RECTANGLE, cx - t * 0.27,
                               cy + t * 0.06, t * 0.54, t * 0.26)
        ombros.adjustments[0] = 0.5
    elif chave == "luminaria":
        _forma_branca(slide, MSO_SHAPE.OVAL, cx - t * 0.19, cy - t * 0.29,
                      t * 0.38, t * 0.38)
        _forma_branca(slide, MSO_SHAPE.RECTANGLE, cx - t * 0.09, cy + t * 0.14,
                      t * 0.18, t * 0.15)
    elif chave == "raio":
        _forma_branca(slide, MSO_SHAPE.LIGHTNING_BOLT, cx - t * 0.17, cy - t * 0.26,
                      t * 0.34, t * 0.52)
    elif chave == "moeda":
        _forma_branca(slide, MSO_SHAPE.DONUT, cx - t * 0.25, cy - t * 0.25,
                      t * 0.50, t * 0.50)
    elif chave == "nuvem":
        _forma_branca(slide, MSO_SHAPE.CLOUD, cx - t * 0.29, cy - t * 0.22,
                      t * 0.58, t * 0.44)
    elif chave == "engrenagem":
        _forma_branca(slide, MSO_SHAPE.GEAR_6, cx - t * 0.27, cy - t * 0.27,
                      t * 0.54, t * 0.54)
    elif chave == "servidor":
        for indice in range(3):
            _forma_branca(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx - t * 0.26,
                          cy - t * 0.27 + indice * t * 0.19, t * 0.52, t * 0.13)
    elif chave == "camera":
        # Corpo branco com a lente recortada em navy: sem o contraste interno o
        # conjunto vira uma mancha branca ilegível no tamanho do emblema.
        corpo = _forma_branca(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx - t * 0.28,
                              cy - t * 0.15, t * 0.46, t * 0.30)
        corpo.adjustments[0] = 0.25
        _forma_branca(slide, MSO_SHAPE.OVAL, cx - t * 0.18, cy - t * 0.06,
                      t * 0.13, t * 0.13, cor=NAVY)
        _forma_branca(slide, MSO_SHAPE.RIGHT_TRIANGLE, cx + t * 0.19, cy - t * 0.10,
                      t * 0.12, t * 0.20)
    elif chave == "wifi":
        # Barras de sinal em vez do leque de arcos: BLOCK_ARC rotacionado se funde
        # numa mancha nesta escala, e barra crescente é lida por qualquer plateia.
        for indice, fracao in enumerate((0.18, 0.30, 0.42, 0.54)):
            largura_barra = t * 0.09
            px = cx - t * 0.28 + indice * t * 0.14
            _forma_branca(slide, MSO_SHAPE.RECTANGLE, px, cy + t * 0.26 - t * fracao,
                          largura_barra, t * fracao)
    elif chave == "fatia":
        _forma_branca(slide, MSO_SHAPE.PIE, cx - t * 0.26, cy - t * 0.26,
                      t * 0.52, t * 0.52)
    else:
        _forma_branca(slide, MSO_SHAPE.OVAL, cx - t * 0.17, cy - t * 0.17,
                      t * 0.34, t * 0.34)


# ── Slide "Pré-viabilidade" (padrão Itabira, p14) ────────────────────────────
def _caixa_valor(slide, x, y, w, h, rotulo_fraco: str, rotulo_forte: str,
                 valor: str) -> None:
    """
    Caixa branca de borda fina preta com rótulo à esquerda e valor à direita.

    É o elemento central do slide de pré-viabilidade do deck original — inclusive a
    borda preta fina, que destoa do resto da paleta mas é o que está no modelo.
    """
    _retangulo(slide, x, y, w, h, BRANCO, borda=RGBColor(0x00, 0x00, 0x00))
    _texto(slide, x + Inches(0.28), y + Inches(0.20), w * 0.45, Inches(0.34),
           rotulo_fraco, tamanho=15, cor=NAVY)
    _texto(slide, x + Inches(0.28), y + Inches(0.56), w * 0.45, Inches(0.34),
           rotulo_forte, tamanho=15, cor=NAVY, negrito=True)
    _texto(slide, x + w * 0.45, y + Inches(0.28), w * 0.55 - Inches(0.28),
           Inches(0.7), valor, tamanho=30, cor=TINTA, negrito=True,
           alinhamento=PP_ALIGN.RIGHT)


def _slide_pre_viabilidade(apresentacao, linha: dict) -> None:
    """
    O slide central do padrão Itabira: arrecadação × despesa, por ponto.mês, e o
    percentual da arrecadação que sobra.

    **Sobre o percentual.** O deck original calcula `(arrecadação − despesa atual) /
    arrecadação` e chama o resultado de "disponível para soluções digitais". O usuário
    determinou em 02/09/2026 replicar esse critério, para manter comparabilidade com as
    peças já entregues. A ressalva técnica fica no rodapé do próprio slide, porque ela
    é material: a despesa declarada é só o custeio de hoje, e a contraprestação da PPP
    — que inclui o CAPEX amortizado, a telegestão e o reinvestimento do ciclo — é
    estruturalmente maior que ela. O indicador que desconta a contraprestação estimada
    é `sobra_percentual`, e aparece no slide de síntese.
    """
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    municipio = str(linha.get("municipio") or "Município")
    uf = str(linha.get("uf") or "")
    exercicio = linha.get("ano_exercicio")
    ano = f" {int(exercicio)}" if exercicio and not pd.isna(exercicio) else ""
    _faixa_titulo(slide, "Pré-viabilidade do projeto",
                  f"{municipio}{'/' + uf if uf else ''}")

    # ── Linha dos pontos ───────────────────────────────────────────────────
    _texto(slide, MARGEM, Inches(1.95), Inches(3.4), Inches(0.34), "PONTOS DE",
           tamanho=15, cor=NAVY)
    _texto(slide, MARGEM, Inches(2.28), Inches(3.4), Inches(0.34),
           "ILUMINAÇÃO PÚBLICA", tamanho=15, cor=NAVY, negrito=True)
    _texto(slide, Inches(5.1), Inches(1.95), Inches(4.4), Inches(0.62),
           f"{_v(linha, 'pontos_ip', 'numero')} pontos", tamanho=30, cor=TINTA)
    _texto(slide, Inches(9.9), Inches(1.95), Inches(3.0), Inches(0.34),
           "ARRECADAÇÃO", tamanho=15, cor=NAVY)
    _texto(slide, Inches(9.9), Inches(2.28), Inches(3.0), Inches(0.34),
           "MÉDIA MENSAL", tamanho=15, cor=NAVY)

    largura_caixa = Inches(8.6)
    despesa = _numero(linha, "despesa_energia_declarada")
    despesa_mes = _numero(linha, "despesa_energia_ponto_mes")

    _caixa_valor(slide, MARGEM, Inches(2.95), largura_caixa, Inches(1.2),
                 "ARRECADAÇÃO", f"CIP{ano}", _v(linha, "cosip_liquida", "reais"))
    _texto(slide, Inches(9.5), Inches(3.22), Inches(2.3), Inches(0.62),
           _v(linha, "cosip_ponto_mes", "reais"), tamanho=28, cor=VERDE_ESCURO,
           negrito=True, alinhamento=PP_ALIGN.RIGHT)
    _texto(slide, Inches(11.9), Inches(3.42), Inches(1.4), Inches(0.34),
           "/ponto.mês", tamanho=13, cor=NAVY)

    rotulo_despesa = "DESPESA DECLARADA" if despesa is not None else "DESPESA COM IP"
    _caixa_valor(slide, MARGEM, Inches(4.35), largura_caixa, Inches(1.2),
                 rotulo_despesa, f"COM ENERGIA{ano}",
                 _v(linha, "despesa_energia_declarada", "reais"))
    _texto(slide, Inches(9.5), Inches(4.62), Inches(2.3), Inches(0.62),
           ranking.formatar_valor(despesa_mes, "reais"), tamanho=28, cor=VERMELHO,
           negrito=True, alinhamento=PP_ALIGN.RIGHT)
    _texto(slide, Inches(11.9), Inches(4.82), Inches(1.4), Inches(0.34),
           "/ponto.mês", tamanho=13, cor=NAVY)

    # ── Percentual disponível ──────────────────────────────────────────────
    disponivel = _disponivel_apos_despesa(linha)
    _glifo(slide, "fatia", MARGEM, Inches(5.75), Inches(0.78))
    if disponivel is not None:
        _texto(slide, MARGEM + Inches(1.05), Inches(5.72), Inches(1.9), Inches(0.72),
               ranking.formatar_valor(disponivel, "percentual"), tamanho=38,
               cor=TINTA, negrito=True)
        _texto(slide, MARGEM + Inches(3.05), Inches(5.82), Inches(8.0), Inches(0.7),
               "Dos valores arrecadados disponíveis para implantação de soluções "
               "digitais", tamanho=15, cor=TINTA, espacamento=1.2)
    else:
        _texto(slide, MARGEM + Inches(1.05), Inches(5.82), Inches(10.5), Inches(0.7),
               "O município não declarou despesa com energia elétrica ao SICONFI no "
               "exercício, e por isso a parcela disponível não é apresentada.",
               tamanho=14, cor=TINTA_FRACA, espacamento=1.2)

    _rodape(slide, _rodape_pre_viabilidade(linha, disponivel))


def _disponivel_apos_despesa(linha: dict) -> Optional[float]:
    """Parcela da arrecadação que sobra depois da despesa declarada — critério do deck."""
    cosip = _numero(linha, "cosip_liquida")
    despesa = _numero(linha, "despesa_energia_declarada")
    if cosip is None or despesa is None or cosip <= 0:
        return None
    return (cosip - despesa) / cosip


def _rodape_pre_viabilidade(linha: dict, disponivel: Optional[float]) -> str:
    """Rodapé com a origem exata da despesa e a ressalva do percentual."""
    partes = []
    origem = str(linha.get("origem_despesa_energia") or "").strip()
    if origem:
        partes.append(f"Despesa declarada na rubrica {origem} (SICONFI, DCA Anexo I-E, "
                      "coluna Liquidada) — é envelope, não gasto exclusivo de "
                      "iluminação pública.")
    if disponivel is not None:
        partes.append("A parcela disponível desconta a despesa de custeio atual, não a "
                      "contraprestação da concessão, que inclui investimento amortizado "
                      "e é maior que o custeio de hoje.")
    partes.append(_fonte_rodape(linha))
    return "  ".join(partes)


# ── Slide "Potenciais impactos" (padrão Itabira, p13) ────────────────────────
def _impactos(linha: dict, prazo_anos: Optional[int]) -> list[tuple[str, str, str]]:
    """
    Os dez itens do slide de impactos: (glifo, manchete, apoio).

    Cinco vêm do que o Hub apurou e mudam com o município; cinco são escopo da
    concessão e não carregam número, porque prometer percentual de redução de
    criminalidade ou de conectividade sem estudo é o tipo de afirmação que não
    sobrevive à primeira pergunta da prefeitura.
    """
    economia_pct = _numero(linha, "economia_percentual")
    tem_economia = (economia_pct is not None
                    and economia_pct >= LIMIAR_ECONOMIA_RELEVANTE)
    populacao = _numero(linha, "populacao")

    calculados: list[tuple[str, str, str]] = [
        ("pessoas", f"{_v(linha, 'populacao', 'numero')} cidadãos"
         if populacao else "População do município", "beneficiados"),
        ("luminaria", f"{_v(linha, 'pontos_ip', 'numero')} pontos",
         "de iluminação pública"),
    ]
    if tem_economia:
        calculados += [
            ("raio", f"{ranking.formatar_valor(economia_pct, 'percentual')} de "
                     "eficientização", "potencial, do consumo de energia"),
            ("moeda", f"{_v(linha, 'economia_ciclo_reais', 'reais_compacto')} de economia",
             f"acumulados em {prazo_anos} anos" if prazo_anos else "acumulados no ciclo"),
            ("nuvem", f"Redução de {_v(linha, 'co2_evitado_t_ano', 'decimal')} "
                      "toneladas de CO₂", "por ano, na matriz elétrica nacional"),
        ]
    else:
        calculados += [
            ("raio", f"{_v(linha, 'perc_led', 'percentual')} do parque já em LED",
             "eficientização em grande parte concluída"),
            ("moeda", f"{_v(linha, 'custo_energia_ano', 'reais_compacto')} de energia",
             "custeio anual do parque como está hoje"),
            ("nuvem", f"{_v(linha, 'potencia_media_w', 'decimal')} W por ponto",
             "potência média instalada"),
        ]

    escopo = [
        ("engrenagem", "100% de telegestão", "da iluminação pública"),
        ("servidor", "Transformação digital", "e integração de serviços públicos"),
        ("camera", "Videomonitoramento", "e apoio à segurança pública"),
        ("wifi", "Conectividade pública", "Wi-Fi em praças e equipamentos"),
        ("fatia", "Custeio pela CIP existente", "sem recursos novos do orçamento"),
    ]
    return calculados + escopo


def _slide_impactos(apresentacao, linha: dict, prazo: Optional[int]) -> None:
    """Grade de dez impactos em duas colunas, no padrão da p13 do deck de Itabira."""
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    municipio = str(linha.get("municipio") or "Município")
    uf = str(linha.get("uf") or "")
    _faixa_titulo(slide, "Potenciais impactos do projeto",
                  f"{municipio}{'/' + uf if uf else ''}")

    prazo_anos = prazo or linha.get("prazo_concessao_anos")
    try:
        prazo_anos = int(float(prazo_anos)) if prazo_anos is not None else None
    except (TypeError, ValueError):
        prazo_anos = None

    itens = _impactos(linha, prazo_anos)
    largura_coluna = (LARGURA - 2 * MARGEM - Inches(0.6)) / 2
    topo, passo = Inches(1.95), Inches(0.93)
    for indice, (glifo, manchete, apoio) in enumerate(itens):
        coluna, fila = indice // 5, indice % 5
        x = MARGEM + coluna * (largura_coluna + Inches(0.6))
        y = topo + fila * passo
        _glifo(slide, glifo, x, y, Inches(0.66))
        _texto(slide, x + Inches(0.92), y + Inches(0.02),
               largura_coluna - Inches(0.92), Inches(0.36), manchete, tamanho=15,
               cor=NAVY, negrito=True)
        _texto(slide, x + Inches(0.92), y + Inches(0.38),
               largura_coluna - Inches(0.92), Inches(0.34), apoio, tamanho=13,
               cor=TINTA)

    _rodape(slide, _fonte_rodape(linha))


# ── Slide "Possíveis escopos" (padrão Itabira, p05) ──────────────────────────
ESCOPOS_CONCESSAO = [
    ("Eficiência energética", VERDE_ESCOPO),
    ("Mobilidade urbana", AZUL),
    ("Segurança pública", NAVY),
    ("Monitoramento ambiental", VERDE_ESCOPO),
    ("Tecnologia e dados", NAVY_ESCURO),
    ("Governança", ROXO),
]


def _cidade_vetorial(slide, x, y, w, h, *, sobre_escuro: bool = False) -> None:
    """
    Silhueta urbana esquemática no centro do slide de escopos.

    O deck original traz uma maquete 3D renderizada. Não há gerador de imagem aqui e
    baixar arte de terceiro para material de cliente traz questão de licença — então o
    lugar dela é ocupado por uma silhueta vetorial, que ninguém confunde com fotografia
    e permanece editável. Se a equipe passar a ilustração oficial, ela entra por cima.
    """
    cor_base = RGBColor(0x0B, 0x4C, 0x86) if sobre_escuro else RGBColor(0xD9, 0xE2, 0xEC)
    cores_predio = ((RGBColor(0x1E, 0x6E, 0xB0), RGBColor(0x14, 0x57, 0x93))
                    if sobre_escuro else (NAVY, AZUL))
    linha_base = y + int(h * 0.72)
    base = _retangulo(slide, x, linha_base, w, Emu(int(h * 0.10)), cor_base,
                      arredondado=True)
    base.adjustments[0] = 0.4
    # (altura, largura) relativas. Larguras diferentes e janelas são o que separa uma
    # silhueta urbana de um gráfico de barras — sem elas o desenho vira um gráfico.
    predios = ((0.30, 0.11), (0.52, 0.09), (0.40, 0.13), (0.66, 0.10),
               (0.34, 0.12), (0.58, 0.09), (0.44, 0.11))
    vao = int(w * 0.018)
    total = sum(int(w * larg) for _, larg in predios) + (len(predios) - 1) * vao
    px = x + int((w - total) / 2)
    cor_janela = BRANCO if sobre_escuro else RGBColor(0xBF, 0xDD, 0xF5)
    for indice, (fracao, largura_rel) in enumerate(predios):
        largura_predio = int(w * largura_rel)
        altura_predio = int(h * 0.72 * fracao)
        py = linha_base - altura_predio
        _retangulo(slide, Emu(px), Emu(py), Emu(largura_predio), Emu(altura_predio),
                   cores_predio[indice % 2])
        _janelas(slide, px, py, largura_predio, altura_predio, cor_janela)
        # Poste de iluminação sobre um prédio baixo — é do que trata a peça.
        if indice == 4:
            _retangulo(slide, Emu(px + int(largura_predio * 0.45)),
                       Emu(py - int(h * 0.16)), Emu(max(int(w * 0.006), 9525)),
                       Emu(int(h * 0.16)), CIANO)
            _forma_branca(slide, MSO_SHAPE.OVAL, px + int(largura_predio * 0.32),
                          py - int(h * 0.19), int(w * 0.030), int(w * 0.030), CIANO)
        px += largura_predio + vao


def _janelas(slide, x, y, w, h, cor) -> None:
    """Grade de janelas de um prédio da silhueta."""
    largura_janela, altura_janela = int(w * 0.16), int(h * 0.055)
    colunas = 3
    passo_x = int(w * 0.28)
    inicio_x = x + int((w - (colunas - 1) * passo_x - largura_janela) / 2)
    passo_y = int(h * 0.14)
    for fila in range(max(int(h * 0.72 / max(passo_y, 1)), 1)):
        py = y + int(h * 0.14) + fila * passo_y
        if py + altura_janela > y + h - int(h * 0.10):
            break
        for coluna in range(colunas):
            _forma_branca(slide, MSO_SHAPE.RECTANGLE, inicio_x + coluna * passo_x, py,
                          largura_janela, altura_janela, cor=cor)


def _slide_escopos(apresentacao) -> None:
    """Os seis escopos possíveis da concessão, em pílulas coloridas ao redor da cidade."""
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    _faixa_titulo(slide, "Possíveis escopos de serviços da concessão")

    _cidade_vetorial(slide, Inches(4.55), Inches(2.30), Inches(4.25), Inches(3.60))

    largura_pilula, altura_pilula = Inches(3.35), Inches(0.72)
    ys = (Inches(2.35), Inches(3.50), Inches(4.65))
    for indice, (rotulo, cor) in enumerate(ESCOPOS_CONCESSAO):
        esquerda = indice % 2 == 0
        x = MARGEM if esquerda else LARGURA - MARGEM - largura_pilula
        y = ys[indice // 2]
        pilula = _retangulo(slide, x, y, largura_pilula, altura_pilula, cor,
                            arredondado=True)
        pilula.adjustments[0] = 0.5
        _texto(slide, x + Inches(0.3), y, largura_pilula - Inches(0.6), altura_pilula,
               rotulo.upper(), tamanho=13, cor=BRANCO, negrito=True,
               alinhamento=PP_ALIGN.CENTER, ancora=MSO_ANCHOR.MIDDLE)
        # Conector até a silhueta, do mesmo jeito que o deck liga pílula e maquete.
        y_meio = y + altura_pilula / 2
        x_inicio = x + largura_pilula if esquerda else Inches(4.55) + Inches(4.25)
        x_fim = Inches(4.55) if esquerda else x
        conector = _retangulo(slide, Emu(int(min(x_inicio, x_fim))),
                              Emu(int(y_meio)), Emu(int(abs(x_fim - x_inicio))),
                              Pt(1.5), cor)
        conector.line.fill.background()

    _texto(slide, MARGEM, Inches(6.15), LARGURA - 2 * MARGEM, Inches(0.7),
           "O escopo efetivo é definido no estudo de viabilidade: a CIP custeia o "
           "serviço de iluminação pública e os sistemas de monitoramento de "
           "logradouros (art. 149-A da Constituição), e cada solução adicional precisa "
           "caber na arrecadação e no interesse do município.",
           tamanho=11, cor=TINTA, espacamento=1.25, alinhamento=PP_ALIGN.CENTER)
    _rodape_marca(slide)


def _slide_tese(apresentacao, linha: dict, prazo: Optional[int],
                contrato: Optional[dict] = None) -> None:
    """
    Slide de fechamento do município: a conta fecha?

    Deliberadamente enxuto. Os números do potencial saíram daqui e foram para o slide
    de impactos (padrão Itabira, p13) — repetir os mesmos seis cartões duas vezes no
    mesmo deck era o efeito colateral óbvio de acrescentar aquele slide. Aqui ficam as
    três coisas que só este slide diz: o custo estimado do serviço, o que o município
    já contrata hoje, e a frase que amarra os dois.
    """
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    municipio = str(linha.get("municipio") or "Município")
    uf = str(linha.get("uf") or "")
    economia_pct = _numero(linha, "economia_percentual")
    tem_economia = economia_pct is not None and economia_pct >= LIMIAR_ECONOMIA_RELEVANTE
    _faixa_titulo(slide, "A conta fecha?", f"{municipio}{'/' + uf if uf else ''}")

    topo = Inches(2.0)
    largura_cartao = (LARGURA - 2 * MARGEM - Inches(0.6)) / 3
    altura_cartao = Inches(1.5)
    blocos = [
        Bloco("Custo estimado do serviço",
              _v(linha, "contraprestacao_mes", "reais_compacto"),
              "contraprestação mensal estimada", NAVY),
        Bloco("Sobra da CIP", _v(linha, "sobra_percentual", "percentual"),
              "da arrecadação, após o serviço", CIANO),
        Bloco("Sobra em reais", _v(linha, "sobra_reais_ano", "reais_compacto"),
              "por ano, para soluções digitais", CIANO),
    ]
    for indice, bloco in enumerate(blocos):
        _cartao(slide, MARGEM + indice * (largura_cartao + Inches(0.3)), topo,
                largura_cartao, altura_cartao, bloco)

    # ── O que o município contrata hoje (PNCP) ─────────────────────────────
    # Entra ANTES da síntese porque é o comparativo direto do custo proposto: sem ele,
    # a contraprestação estimada fica sem referência na cabeça de quem assiste.
    y_atual = topo + altura_cartao + Inches(0.35)
    if contrato:
        _faixa_contrato(slide, contrato, y_atual)
        y_atual += Inches(1.2)

    # Frase-síntese, condicionada ao que os números realmente dizem.
    sobra = _numero(linha, "sobra_percentual")
    faixa = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGEM, y_atual,
                                   LARGURA - 2 * MARGEM, Inches(1.5))
    faixa.adjustments[0] = 0.12
    _degrade(faixa, NAVY, AZUL)
    faixa.text_frame.text = ""
    _texto(slide, MARGEM + Inches(0.4), y_atual + Inches(0.24),
           LARGURA - 2 * MARGEM - Inches(0.8), Inches(1.1),
           _sintese(sobra, economia_pct, tem_economia), tamanho=13, cor=BRANCO,
           espacamento=1.25)
    _rodape(slide, _fonte_rodape(linha))


def _faixa_contrato(slide, contrato: dict, y) -> None:
    """
    Faixa com o contrato de iluminação que o município já tem publicado no PNCP.

    Duas cautelas, ambas visíveis no slide:

      - o valor citado é o **equivalente mensal pelo prazo contratual**, não o valor
        global — comparar valor global com contraprestação mensal erraria a ordem de
        grandeza em uma casa ou mais;
      - o texto diz que a fonte é o PNCP e mostra a vigência, para que a prefeitura
        confira. A busca do portal é textual e pode trazer contrato de objeto vizinho.
    """
    largura = LARGURA - 2 * MARGEM
    _retangulo(slide, MARGEM, y, largura, Inches(0.98), BRANCO,
               borda=CINZA_CLARO, arredondado=True)
    _retangulo(slide, MARGEM, y, Inches(0.05), Inches(0.98), AMBAR)
    _texto(slide, MARGEM + Inches(0.22), y + Inches(0.14), largura - Inches(0.44),
           Inches(0.28), "O QUE O MUNICÍPIO CONTRATA HOJE — PNCP", tamanho=9,
           cor=TINTA_FRACA, negrito=True)

    mensal = contrato.get("valor_mensal_equivalente")
    global_ = contrato.get("valor_global")
    inicio = _data_br(contrato.get("data_inicio_vigencia"))
    fim = _data_br(contrato.get("data_fim_vigencia"))
    vigencia = f"vigência {inicio} a {fim}" if inicio and fim else "vigência não informada"

    if mensal:
        principal_txt = (f"{ranking.formatar_valor(mensal, 'reais_compacto')}/mês "
                         "equivalente pelo prazo contratual")
    elif global_:
        principal_txt = (f"{ranking.formatar_valor(global_, 'reais_compacto')} "
                         "de valor global")
    else:
        principal_txt = "contrato publicado sem valor informado"

    _texto(slide, MARGEM + Inches(0.22), y + Inches(0.42), largura - Inches(0.44),
           Inches(0.46), f"{principal_txt} · {vigencia}.", tamanho=12, cor=TINTA,
           espacamento=1.15)


def _sintese(sobra: Optional[float], economia: Optional[float],
             tem_economia: bool = True) -> str:
    """
    Frase de fechamento — condicionada ao número, nunca genérica.

    Sobra negativa não vira silêncio nem promessa: vira a informação de que a CIP
    atual não banca o serviço sozinha, que é um fato com o qual a prefeitura precisa
    lidar (revisão da lei de CIP, aporte ou escopo menor).
    """
    if not tem_economia:
        base = ("O parque já está majoritariamente modernizado, de modo que o ganho "
                "da concessão não está na troca de luminárias, e sim na operação e "
                "manutenção permanentes, na telegestão, na expansão da rede e no "
                "atendimento com prazo contratado. ")
        if sobra is not None and sobra >= 0:
            return base + ("A CIP hoje arrecadada comporta esse serviço, com folga "
                           f"estimada de {ranking.formatar_valor(sobra, 'percentual')}.")
        return base + ("Nas premissas adotadas, a CIP hoje arrecadada não comporta "
                       "esse serviço sem revisão da lei de custeio ou ajuste de escopo.")

    if sobra is None:
        return ("A arrecadação de CIP e o parque levantado permitem estruturar a "
                "modernização; o dimensionamento final depende de cadastro validado "
                "em campo.")
    if sobra >= 0.15:
        return ("A CIP hoje arrecadada é suficiente para custear a modernização e a "
                f"operação do parque, com folga estimada de "
                f"{ranking.formatar_valor(sobra, 'percentual')} da arrecadação — sem "
                "recursos novos do orçamento municipal.")
    if sobra >= 0:
        return ("A CIP hoje arrecadada cobre a modernização e a operação do parque, "
                f"com folga estreita de {ranking.formatar_valor(sobra, 'percentual')}. "
                "O equilíbrio depende de cadastro validado e de escopo bem delimitado.")
    return ("A CIP hoje arrecadada não cobre sozinha a modernização e a operação nas "
            "premissas adotadas. Viabilizar o projeto exige revisar a lei de custeio, "
            "ajustar o escopo ou prever contrapartida orçamentária.")


def _fonte_rodape(linha: dict) -> str:
    partes = ["Fontes: SICONFI/Tesouro Nacional (DCA), BDGD/ANEEL, IBGE."]
    origem = str(linha.get("origem_pontos") or "")
    if "estimad" in origem.lower():
        partes.append("Parque estimado por população — município sem base da "
                      "distribuidora.")
    defasagem = linha.get("defasagem_anos")
    try:
        if defasagem is not None and not pd.isna(defasagem) and float(defasagem) >= 2:
            partes.append(f"Base da distribuidora {int(float(defasagem))} anos "
                          "anterior ao exercício da receita.")
    except (TypeError, ValueError):
        pass
    partes.append("Valores estimados para triagem; não substituem estudo de "
                  "viabilidade com cadastro validado em campo.")
    return "  ".join(partes)


def _slide_metodologia(apresentacao, premissas: dict) -> None:
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    _faixa_titulo(slide, "Metodologia e fontes",
                  "De onde vem cada número desta apresentação")

    itens = [
        ("Arrecadação de CIP",
         "Declaração do município ao SICONFI/Tesouro Nacional, DCA Anexo I-C, "
         "receita líquida do exercício."),
        ("Despesa com energia elétrica",
         "SICONFI, DCA Anexo I-E, subfunção 25.752 (ou 25.751, onde o município declara "
         "assim), coluna Liquidada. É a energia de todas as unidades do município — "
         "iluminação, prédios e escolas —, e por isso um teto para o gasto com "
         "iluminação, não o gasto exato. Onde o município não declara a função 25, o "
         "indicador não é apresentado."),
        ("Parque de iluminação",
         "Base de Dados Geográfica da Distribuidora (BDGD/ANEEL), entidade PIP, pontos "
         "ativos. Onde não há base disponível, o parque é estimado pela população e "
         "isso vem sinalizado no slide."),
        ("Economia de energia",
         "Redução proporcional da potência média do parque até a potência de referência "
         "adotada, aplicada ao consumo. O acumulado do ciclo é nominal, com reajuste "
         "anual, pela soma da série geométrica."),
        ("Contratação atual",
         "Portal Nacional de Contratações Públicas (PNCP). O valor citado é o "
         "equivalente mensal pelo prazo do contrato, obtido do valor global publicado. "
         "A busca do portal é textual: pode não trazer contrato anterior à Lei "
         "14.133/2021 nem distinguir objetos vizinhos."),
        ("Localização",
         "Malhas territoriais do IBGE."),
    ]
    # O passo é apertado de propósito: com espaçamento folgado o sexto item cai por
    # baixo da tarja de premissas e some do slide sem que nada acuse.
    y = Inches(1.85)
    for titulo, corpo in itens:
        _texto(slide, MARGEM, y, Inches(3.1), Inches(0.4), titulo, tamanho=11.5,
               cor=NAVY, negrito=True)
        _texto(slide, MARGEM + Inches(3.3), y, LARGURA - MARGEM * 2 - Inches(3.3),
               Inches(0.8), corpo, tamanho=10, cor=TINTA, espacamento=1.15)
        y += Inches(0.68)

    premissas_texto = " · ".join(
        f"{chave}: {valor}" for chave, valor in premissas.items() if valor is not None)
    if premissas_texto:
        _retangulo(slide, MARGEM, ALTURA - Inches(1.55), LARGURA - 2 * MARGEM,
                   Inches(0.55), CINZA_CLARO, arredondado=True)
        _texto(slide, MARGEM + Inches(0.25), ALTURA - Inches(1.40),
               LARGURA - 2 * MARGEM - Inches(0.5), Inches(0.4),
               f"Premissas da triagem — {premissas_texto}", tamanho=9.5, cor=TINTA)
    _rodape(slide, "Triagem de oportunidades de modernização de iluminação pública a "
                   "partir de bases públicas federais.")


def _icone(slide, x, y, tamanho, indice: int) -> None:
    """
    Marca visual do cartão de solução.

    Ícones geométricos, desenhados com as formas nativas do PowerPoint: mantêm a
    apresentação 100% vetorial e editável, sem imagem de terceiro — que traria questão
    de licença num material que vai para o cliente.
    """
    circulo = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, tamanho, tamanho)
    circulo.fill.solid()
    circulo.fill.fore_color.rgb = NAVY
    circulo.line.fill.background()
    circulo.shadow.inherit = False

    centro_x, centro_y = x + tamanho // 2, y + tamanho // 2
    miolo = int(tamanho * 0.34)
    formas = [MSO_SHAPE.OVAL, MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.DIAMOND,
              MSO_SHAPE.PENTAGON, MSO_SHAPE.OVAL, MSO_SHAPE.RECTANGLE]
    interno = slide.shapes.add_shape(
        formas[indice % len(formas)], Emu(int(centro_x - miolo / 2)),
        Emu(int(centro_y - miolo / 2)), Emu(miolo), Emu(miolo))
    interno.fill.solid()
    interno.fill.fore_color.rgb = CIANO
    interno.line.fill.background()
    interno.shadow.inherit = False


def _slide_contexto(apresentacao, linhas: list[dict]) -> None:
    """
    Abre a apresentação de cliente situando de onde vêm os números.

    Sem isso a primeira coisa que a prefeitura vê é um número sobre a própria cidade,
    sem saber quem apurou nem com base em quê — e a conversa começa pela desconfiança.
    """
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    _faixa_titulo(slide, "De onde vem esta análise",
                  "Dados públicos, cruzados e verificáveis")

    itens = [
        ("Arrecadação", "O que o município declarou de CIP ao Tesouro Nacional."),
        ("Parque", "Os pontos de iluminação registrados pela distribuidora na "
                   "base da ANEEL."),
        ("Despesa", "O que o município declarou gastar com energia elétrica."),
        ("Potencial", "O que a modernização economiza, e o que a arrecadação "
                      "existente é capaz de custear."),
    ]
    largura = (LARGURA - 2 * MARGEM - Inches(0.9)) / 4
    for indice, (titulo, corpo) in enumerate(itens):
        x = MARGEM + indice * (largura + Inches(0.3))
        _retangulo(slide, x, Inches(1.95), largura, Inches(2.3), BRANCO,
                   borda=CINZA_CLARO, arredondado=True)
        _retangulo(slide, x, Inches(1.95), largura, Inches(0.06), CIANO)
        _texto(slide, x + Inches(0.24), Inches(2.25), largura - Inches(0.48),
               Inches(0.4), titulo, tamanho=15, cor=NAVY, negrito=True)
        _texto(slide, x + Inches(0.24), Inches(2.73), largura - Inches(0.48),
               Inches(1.4), corpo, tamanho=11, cor=TINTA, espacamento=1.25)

    _retangulo(slide, MARGEM, Inches(4.65), LARGURA - 2 * MARGEM, Inches(1.35), NAVY,
               arredondado=True)
    _texto(slide, MARGEM + Inches(0.35), Inches(4.92), LARGURA - 2 * MARGEM - Inches(0.7),
           Inches(0.95),
           "Nenhum dado desta apresentação foi coletado do município: tudo vem de "
           "bases públicas federais, e pode ser conferido nas fontes indicadas ao "
           "final. Os valores são de triagem — o dimensionamento definitivo exige "
           "cadastro validado em campo.",
           tamanho=13, cor=BRANCO, espacamento=1.25)
    _rodape(slide, "Fontes: SICONFI/Tesouro Nacional, BDGD/ANEEL, IBGE, PNCP.")


def _slide_solucoes(apresentacao, linha: dict) -> None:
    """
    O que a sobra da CIP pode custear além da iluminação.

    Só é gerado quando há sobra positiva: oferecer câmera e Wi-Fi a município cuja
    arrecadação não cobre nem a operação da iluminação seria vender o que a conta não
    paga — exatamente o tipo de promessa que não sobrevive à primeira due diligence.
    """
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    municipio = str(linha.get("municipio") or "Município")
    uf = str(linha.get("uf") or "")
    _faixa_titulo(slide, f"{municipio}{'/' + uf if uf else ''}",
                  "O que a sobra da CIP pode custear além da iluminação")

    sobra_reais = _numero(linha, "sobra_reais_ano")
    if sobra_reais:
        _texto(slide, MARGEM, Inches(1.85), LARGURA - 2 * MARGEM, Inches(0.4),
               f"Com a modernização, sobram cerca de "
               f"{ranking.formatar_valor(sobra_reais, 'reais_compacto')} por ano da "
               "arrecadação já existente. Esse saldo pode financiar, na mesma "
               "concessão:", tamanho=12, cor=TINTA)

    topo = Inches(2.35)
    largura = (LARGURA - 2 * MARGEM - Inches(0.6)) / 3
    altura = Inches(1.70)
    for indice, (titulo, descricao) in enumerate(SOLUCOES_DIGITAIS):
        coluna, fila = indice % 3, indice // 3
        x = MARGEM + coluna * (largura + Inches(0.3))
        y = topo + fila * (altura + Inches(0.28))
        _retangulo(slide, x, y, largura, altura, BRANCO, borda=CINZA_CLARO,
                   arredondado=True)
        _icone(slide, x + Inches(0.24), y + Inches(0.24), Inches(0.52), indice)
        _texto(slide, x + Inches(0.92), y + Inches(0.34), largura - Inches(1.16),
               Inches(0.35), titulo, tamanho=13, cor=NAVY, negrito=True)
        _texto(slide, x + Inches(0.24), y + Inches(0.92), largura - Inches(0.48),
               Inches(0.8), descricao, tamanho=10, cor=TINTA_FRACA, espacamento=1.2)

    _rodape(slide, "Escopo e dimensionamento das soluções digitais são definidos na "
                   "proposta, conforme a prioridade do município e o saldo efetivo "
                   "da arrecadação.  " + _fonte_rodape(linha))


# ── Modo comercial ───────────────────────────────────────────────────────────
def _slide_panorama_carteira(apresentacao, linhas: list[dict], titulo: str) -> None:
    """Abre a apresentação comercial com o tamanho do funil."""
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    _faixa_titulo(slide, "Panorama da carteira", titulo)

    total_pontos = sum(_numero(l, "pontos_ip") or 0 for l in linhas)
    total_contra = sum(_numero(l, "contraprestacao_mes") or 0 for l in linhas)
    quentes = sum(1 for l in linhas if str(l.get("temperatura")) == "Quente")
    sem_ppp = sum(1 for l in linhas if not bool(l.get("tem_ppp")))

    blocos = [
        Bloco("Municípios na carteira", ranking.formatar_valor(len(linhas), "numero"),
              f"{quentes} classificados como quentes", CIANO),
        Bloco("Pontos de IP somados", ranking.formatar_valor(total_pontos, "numero"),
              "parque total dos leads", CIANO),
        Bloco("Contraprestação potencial",
              ranking.formatar_valor(total_contra, "reais_compacto") + "/mês",
              "soma dos contratos estimados", NAVY),
        Bloco("Sem PPP contratada", ranking.formatar_valor(sem_ppp, "numero"),
              "ativo ainda disponível", NAVY),
    ]
    largura = (LARGURA - 2 * MARGEM - Inches(0.9)) / 4
    for indice, bloco in enumerate(blocos):
        _cartao(slide, MARGEM + indice * (largura + Inches(0.3)), Inches(1.95),
                largura, Inches(1.5), bloco)

    _texto(slide, MARGEM, Inches(3.75), LARGURA - 2 * MARGEM, Inches(0.4),
           "Como o score é montado", tamanho=15, cor=NAVY, negrito=True)
    explicacao = [
        ("Sobra da CIP — 40%", "Quanto da arrecadação sobra depois de paga a "
                               "contraprestação. Folga grande é projeto que fecha sem "
                               "negociação difícil."),
        ("Tamanho do contrato — 35%", "Pontos × custo por ponto. Receita potencial do "
                                      "contrato."),
        ("Parque legado sem PPP — 25%", "Potência média alta e pouco LED, com "
                                        "concessão ainda não contratada. Município que "
                                        "já tem PPP zera este eixo."),
    ]
    y = Inches(4.25)
    for titulo_item, corpo in explicacao:
        _texto(slide, MARGEM, y, Inches(3.4), Inches(0.4), titulo_item, tamanho=11,
               cor=CIANO, negrito=True)
        _texto(slide, MARGEM + Inches(3.6), y, LARGURA - 2 * MARGEM - Inches(3.6),
               Inches(0.7), corpo, tamanho=11, cor=TINTA, espacamento=1.2)
        y += Inches(0.72)

    _rodape(slide, "Score relativo a esta carteira: o mesmo município pode ser quente "
                   "numa lista regional e morno numa lista nacional. Municípios com "
                   "dado suspeito têm o score reduzido pela metade.")


def _slide_lead(apresentacao, linha: dict, posicao: int, contato: Optional[dict],
                contrato: Optional[dict]) -> None:
    """Ficha de abordagem de um lead, para o time comercial."""
    slide = _slide_vazio(apresentacao)
    _fundo_padrao(slide)
    municipio = str(linha.get("municipio") or "Município")
    uf = str(linha.get("uf") or "")
    temperatura = str(linha.get("temperatura") or "")
    _faixa_titulo(slide, f"{posicao}. {municipio}{'/' + uf if uf else ''}",
                  f"Lead {temperatura.lower()} · score "
                  f"{ranking.formatar_valor(linha.get('score_comercial'), 'decimal')}"
                  if temperatura else "Ficha de abordagem")

    largura = (LARGURA - 2 * MARGEM - Inches(0.9)) / 4
    blocos = [
        Bloco("Pontos de IP", _v(linha, "pontos_ip", "numero"),
              str(linha.get("origem_pontos") or "BDGD"), CIANO),
        Bloco("Contraprestação", _v(linha, "contraprestacao_mes", "reais_compacto"),
              "potencial por mês", CIANO),
        Bloco("Sobra da CIP", _v(linha, "sobra_percentual", "percentual"),
              "da arrecadação", NAVY),
        Bloco("Parque em LED", _v(linha, "perc_led", "percentual"),
              _v(linha, "potencia_media_w", "decimal") + " W médios", NAVY),
    ]
    for indice, bloco in enumerate(blocos):
        _cartao(slide, MARGEM + indice * (largura + Inches(0.3)), Inches(1.95),
                largura, Inches(1.45), bloco)

    # ── Situação contratual ────────────────────────────────────────────────
    y = Inches(3.60)
    _texto(slide, MARGEM, y, Inches(6), Inches(0.32), "SITUAÇÃO CONTRATUAL",
           tamanho=9, cor=TINTA_FRACA, negrito=True)
    if linha.get("tem_ppp"):
        situacao = (f"PPP já contratada com {linha.get('concessionaria_ppp') or 'concessionária'}"
                    f"{' desde ' + str(int(float(linha['ano_ppp']))) if linha.get('ano_ppp') else ''}. "
                    "Ativo tomado — abordagem só faz sentido para escopo complementar.")
    elif contrato:
        mensal = contrato.get("valor_mensal_equivalente")
        situacao = ("Sem PPP. Contrato de manutenção publicado no PNCP"
                    + (f", equivalente a {ranking.formatar_valor(mensal, 'reais_compacto')}/mês"
                       if mensal else "")
                    + (f", vigente até {_data_br(contrato.get('data_fim_vigencia'))}"
                       if contrato.get("data_fim_vigencia") else "") + ".")
    else:
        situacao = ("Sem PPP contratada e sem contrato de manutenção localizado no "
                    "PNCP — verificar se a operação é feita por equipe própria.")
    _texto(slide, MARGEM, y + Inches(0.3), LARGURA - 2 * MARGEM, Inches(0.6),
           situacao, tamanho=12, cor=TINTA, espacamento=1.2)

    # ── Contato institucional ──────────────────────────────────────────────
    y_contato = Inches(4.70)
    _retangulo(slide, MARGEM, y_contato, LARGURA - 2 * MARGEM, Inches(1.55), BRANCO,
               borda=CINZA_CLARO, arredondado=True)
    _retangulo(slide, MARGEM, y_contato, Inches(0.05), Inches(1.55), CIANO)
    _texto(slide, MARGEM + Inches(0.25), y_contato + Inches(0.16),
           Inches(6), Inches(0.3), "CANAL INSTITUCIONAL", tamanho=9,
           cor=TINTA_FRACA, negrito=True)

    campos = [
        ("Site", (contato or {}).get("site")),
        ("E-mail", (contato or {}).get("email")),
        ("Telefone", (contato or {}).get("telefone")),
    ]
    x_campo = MARGEM + Inches(0.25)
    largura_campo = (LARGURA - 2 * MARGEM - Inches(0.5)) / 3
    for rotulo_campo, valor in campos:
        _texto(slide, x_campo, y_contato + Inches(0.52), largura_campo, Inches(0.28),
               rotulo_campo, tamanho=9, cor=TINTA_FRACA, negrito=True)
        _texto(slide, x_campo, y_contato + Inches(0.78), largura_campo, Inches(0.5),
               str(valor) if valor else "— preencher —", tamanho=11,
               cor=TINTA if valor else TINTA_FRACA, espacamento=1.15)
        x_campo += largura_campo

    if contato and contato.get("site"):
        _texto(slide, MARGEM + Inches(0.25), y_contato + Inches(1.26),
               LARGURA - 2 * MARGEM - Inches(0.5), Inches(0.26),
               f"Coletado do portal oficial em {contato.get('coletado_em') or '—'} — "
               "confirme antes de acionar.", tamanho=8, cor=TINTA_FRACA)

    _rodape(slide, "Uso interno — não distribuir ao município. " + _fonte_rodape(linha))


# ── Entrada pública ──────────────────────────────────────────────────────────
PUBLICO_CLIENTE = "cliente"
PUBLICO_COMERCIAL = "comercial"
PUBLICOS = [PUBLICO_CLIENTE, PUBLICO_COMERCIAL]


def gerar(
    painel: pd.DataFrame,
    titulo: str = "Modernização da Iluminação Pública",
    subtitulo: str = ("Panorama da arrecadação de CIP, do parque instalado e do "
                      "potencial de modernização."),
    premissas: Optional[dict] = None,
    contratos_pncp: Optional[dict] = None,
    publico: str = PUBLICO_CLIENTE,
    contatos: Optional[dict] = None,
    incluir_solucoes: bool = True,
) -> bytes:
    """
    Monta a apresentação .pptx dos municípios selecionados.

    Dois públicos, com conteúdos deliberadamente diferentes:

      - **cliente** (padrão): vai para a prefeitura, no padrão visual do deck de
        Itabira. Capa, contexto das fontes, escopos possíveis da concessão e, por
        município, panorama → pré-viabilidade → potenciais impactos → a conta fecha? →
        soluções digitais custeáveis pela sobra. Fecha com a metodologia. Não expõe
        score, ranking, contato nem comparação entre municípios — a cidade não quer
        saber que é a sétima da lista.
      - **comercial**: uso interno. Panorama do funil, fichas de abordagem ordenadas
        por score, situação contratual e canal institucional. Cada slide é carimbado
        como uso interno, porque é material que não pode vazar para o município.

    Args:
        painel: linhas do painel de indicadores (uma por município).
        publico: `PUBLICO_CLIENTE` ou `PUBLICO_COMERCIAL`.
        contratos_pncp: código IBGE → contrato (saída de `contratos_ip.principal`).
        contatos: código IBGE → contato institucional (`contato_municipio.Contato`
            como dict). Usado apenas no modo comercial.
        incluir_solucoes: no modo cliente, acrescenta o slide de soluções digitais
            para quem tem sobra positiva.

    Returns:
        Bytes do arquivo .pptx.
    """
    if publico not in PUBLICOS:
        raise ValueError(f"Público desconhecido: {publico!r}. Use um de {PUBLICOS}.")

    apresentacao = Presentation()
    apresentacao.slide_width = LARGURA
    apresentacao.slide_height = ALTURA

    linhas = [] if painel is None or painel.empty else painel.to_dict("records")
    contratos = contratos_pncp or {}
    fichas = contatos or {}

    prazo = None
    if linhas:
        try:
            prazo = int(float(linhas[0].get("prazo_concessao_anos")))
        except (TypeError, ValueError):
            prazo = None

    # A capa mostra o mapa do município quando a peça é de um só — que é o caso do
    # deck de referência. Em carteira com vários, mapa nenhum representaria o conjunto.
    linha_capa = linhas[0] if len(linhas) == 1 else None

    if publico == PUBLICO_COMERCIAL:
        _capa(apresentacao, titulo, subtitulo, len(linhas), linha_capa)
        if linhas:
            _slide_panorama_carteira(apresentacao, linhas, subtitulo)
        for posicao, linha in enumerate(linhas, start=1):
            chave = str(linha.get("codigo_municipio") or "")
            _slide_lead(apresentacao, linha, posicao, fichas.get(chave),
                        contratos.get(chave))
    else:
        _capa(apresentacao, titulo, subtitulo, len(linhas), linha_capa)
        _slide_contexto(apresentacao, linhas)
        _slide_escopos(apresentacao)
        for linha in linhas:
            chave = str(linha.get("codigo_municipio") or "")
            _slide_panorama(apresentacao, linha)
            _slide_pre_viabilidade(apresentacao, linha)
            _slide_impactos(apresentacao, linha, prazo)
            _slide_tese(apresentacao, linha, prazo, contratos.get(chave))
            # Soluções digitais só onde a conta fecha: oferecer câmera e Wi-Fi a
            # município cuja CIP não cobre nem a operação seria vender o que a
            # arrecadação não paga.
            sobra = _numero(linha, "sobra_percentual")
            if incluir_solucoes and sobra is not None and sobra > 0:
                _slide_solucoes(apresentacao, linha)

    _slide_metodologia(apresentacao, premissas or {})

    buffer = io.BytesIO()
    apresentacao.save(buffer)
    return buffer.getvalue()


__all__ = ["gerar", "desenhar_mapa", "NAVY", "AZUL", "CIANO", "VERDE", "VERMELHO",
           "PUBLICO_CLIENTE", "PUBLICO_COMERCIAL", "PUBLICOS",
           "SOLUCOES_DIGITAIS", "ESCOPOS_CONCESSAO"]
